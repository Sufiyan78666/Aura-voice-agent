"""
rag_tool.py — Multimodal RAG for Voice Agent
Supports:
  Text  : PDF, TXT, MD, DOCX, DOC, PPTX, PPT, XLSX, XLS, CSV
  Images: Standalone (PNG/JPG/etc) + embedded in DOCX, PPTX
          (PDF image extraction disabled — text PDFs don't need it)
Embedding : nomic-embed-text (local Ollama) — max 2048 tokens, safe chunk=512
Vision    : llava (local Ollama) via direct API — no llama-index dependency
Query LLM : gemma4:31b-cloud (cloud)

OCR       : ocrmypdf + Tesseract — auto-applied to scanned (image-only) PDFs.
            Install: pip install ocrmypdf  +  Tesseract on PATH
            Skipped gracefully if ocrmypdf is not installed.
"""

from __future__ import annotations

import hashlib
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

# ─────────────────────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────────────────────
RAG_DOCS_DIR     = os.environ.get("RAG_DOCS_DIR",      "./rag_docs")
RAG_INDEX_DIR    = os.environ.get("RAG_INDEX_DIR",     "./rag_index")
EMBED_MODEL      = os.environ.get("RAG_EMBED_MODEL",   "nomic-embed-text")
EMBED_BACKEND    = os.environ.get("RAG_EMBED_BACKEND", "ollama")
VISION_MODEL     = os.environ.get("RAG_VISION_MODEL",  "llava")
OLLAMA_LOCAL_URL = "http://localhost:11434"
CHROMA_PERSIST_DIR = os.environ.get("CHROMA_PERSIST_DIR", RAG_INDEX_DIR)
CHROMA_COLLECTION  = os.environ.get("CHROMA_COLLECTION", "voice_agent_rag")

# nomic-embed-text hard limit is 2048 tokens (~1500 chars).
# chunk_size=512 tokens is safe and gives good retrieval granularity.
SAFE_CHUNK_SIZE  = int(os.environ.get("RAG_CHUNK_SIZE", "512"))
MAX_NODE_CHARS   = 1800   # hard truncation safety net
MIN_IMAGE_BYTES  = int(os.environ.get("RAG_MIN_IMAGE_BYTES", "5000"))  # skip images smaller than this

TEXT_EXTS  = {".pdf", ".txt", ".md", ".docx", ".doc", ".pptx", ".ppt"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"}
EXCEL_EXTS = {".xlsx", ".xls", ".csv"}

_index_cache = None
_HASH_FILE   = ".rag_docs_hash"



# ─────────────────────────────────────────────────────────────
# TEXT QUALITY FILTERING
# ─────────────────────────────────────────────────────────────
def _clean_text(text: str) -> str:
    """Remove control chars and normalize whitespace for safer indexing."""
    if not text:
        return ""
    cleaned = []
    for ch in text:
        code = ord(ch)
        if ch in ("\n", "\t"):
            cleaned.append(ch)
        elif 32 <= code <= 126:
            cleaned.append(ch)
    cleaned_text = "".join(cleaned)
    return "\n".join(line.rstrip() for line in cleaned_text.splitlines()).strip()


def _is_text_quality_ok(
    text: str,
    min_chars: int = 30,
    min_alnum_ratio: float = 0.2,
    min_words: int = 3,
) -> bool:
    stripped = (text or "").strip()
    if not stripped:
        return False
    if len(stripped) < min_chars:
        return len(stripped.split()) >= min_words
    non_space = [c for c in stripped if not c.isspace()]
    if not non_space:
        return False
    alnum = sum(c.isalnum() for c in non_space)
    return (alnum / len(non_space)) >= min_alnum_ratio


def _filter_documents(docs: list["Document"]) -> list["Document"]:
    """Clean and drop documents that look like OCR garbage or binary noise."""
    from llama_index.core import Document

    kept = []
    skipped = 0
    for doc in docs:
        cleaned = _clean_text(getattr(doc, "text", ""))
        if not _is_text_quality_ok(cleaned):
            skipped += 1
            continue
        metadata = getattr(doc, "metadata", None) or {}
        kept.append(Document(text=cleaned, metadata=metadata))
    if skipped:
        print(f"⚠️  Skipped {skipped} low-quality text doc(s).")
    return kept


# ─────────────────────────────────────────────────────────────
# CHANGE DETECTION
# ─────────────────────────────────────────────────────────────
def _compute_docs_hash(docs_dir: str) -> str:
    docs_path = Path(docs_dir)
    if not docs_path.exists():
        return ""
    entries = []
    for f in sorted(docs_path.rglob("*")):
        if f.is_file():
            stat = f.stat()
            entries.append(f"{f.name}|{stat.st_size}|{stat.st_mtime}")
    return hashlib.md5("\n".join(entries).encode()).hexdigest()


def _load_saved_hash(index_dir: str) -> str:
    try:
        return (Path(index_dir) / _HASH_FILE).read_text().strip()
    except Exception:
        return ""


def _save_hash(docs_dir: str, index_dir: str):
    hash_path = Path(index_dir) / _HASH_FILE
    hash_path.parent.mkdir(parents=True, exist_ok=True)
    hash_path.write_text(_compute_docs_hash(docs_dir))


def docs_changed(docs_dir: str, index_dir: str) -> bool:
    return _compute_docs_hash(docs_dir) != _load_saved_hash(index_dir)


# ─────────────────────────────────────────────────────────────
# OCR — SCANNED PDF DETECTION & AUTO-CORRECTION
# ─────────────────────────────────────────────────────────────
def _is_scanned_pdf(pdf_path: str, text_threshold: int = 50) -> bool:
    """
    Return True if the PDF has little/no extractable text (i.e. it's a scan).
    Uses pdfminer or pypdf to attempt text extraction on the first few pages.
    Falls back to True (assume scanned) if neither library is available.

    text_threshold: minimum number of characters across sampled pages to be
                    considered a 'text PDF'. 50 chars is very conservative.
    """
    # Try pdfminer first (more reliable for text detection)
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        sample = pdfminer_extract(pdf_path, maxpages=3)
        char_count = len((sample or "").strip())
        print(f"   🔍 Text chars found (pdfminer, 3 pages): {char_count}")
        return char_count < text_threshold
    except ImportError:
        pass

    # Try pypdf as fallback
    try:
        from pypdf import PdfReader
        reader = PdfReader(pdf_path)
        text = ""
        for page in reader.pages[:3]:
            text += page.extract_text() or ""
        char_count = len(text.strip())
        print(f"   🔍 Text chars found (pypdf, 3 pages): {char_count}")
        return char_count < text_threshold
    except ImportError:
        pass

    # Can't detect — assume scanned to be safe
    print("   ⚠️  Neither pdfminer nor pypdf installed — assuming scanned PDF.")
    return True


def _ocr_pdf(pdf_path: str, output_path: str) -> bool:
    """
    Run ocrmypdf on a scanned PDF to produce a searchable text PDF.
    Returns True on success, False if ocrmypdf is unavailable or fails.

    ocrmypdf flags used:
      --skip-text        : skip pages that already have a text layer
      --optimize 0       : no image compression (preserve quality)
      --output-type pdf  : standard PDF output (not pdfa)
      -l eng             : English language (change via RAG_OCR_LANG env var)
    """
    try:
        import ocrmypdf  # noqa: F401 — just check it's installed
    except ImportError:
        print("   ⚠️  ocrmypdf not installed — scanned PDF will be skipped.")
        print("       Install with: pip install ocrmypdf")
        print("       Also install Tesseract: https://github.com/UB-Mannheim/tesseract/wiki")
        return False

    ocr_lang = os.environ.get("RAG_OCR_LANG", "eng")
    cmd = [
        "ocrmypdf",
        "--skip-text",       # don't re-OCR pages that already have text
        "--optimize", "0",   # preserve image quality
        "--output-type", "pdf",
        "-l", ocr_lang,
        pdf_path,
        output_path,
    ]

    print(f"   🔠 Running OCR on {Path(pdf_path).name} (lang={ocr_lang})...")
    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=300,  # 5-minute timeout per PDF
        )
        if result.returncode == 0:
            print(f"   ✅ OCR complete → {Path(output_path).name}")
            return True
        else:
            print(f"   ❌ ocrmypdf failed (exit {result.returncode}):")
            if result.stderr:
                # Print last 5 lines of stderr (most relevant)
                for line in result.stderr.strip().splitlines()[-5:]:
                    print(f"      {line}")
            return False
    except FileNotFoundError:
        print("   ❌ 'ocrmypdf' command not found — is it on PATH?")
        return False
    except subprocess.TimeoutExpired:
        print(f"   ❌ OCR timed out for {Path(pdf_path).name}")
        return False


def _preprocess_pdfs(docs_dir: str) -> dict[str, str]:
    """
    Scan all PDFs in docs_dir. For any that look like scans, run OCR and
    write a searchable PDF to a temp subfolder.

    Returns a mapping:  original_path → ocr_path  (only for OCR'd files).
    Non-scanned PDFs are not in the mapping and are used as-is.
    """
    docs_path  = Path(docs_dir)
    ocr_dir    = docs_path / "_ocr_cache"
    ocr_dir.mkdir(exist_ok=True)
    ocr_map: dict[str, str] = {}

    pdf_files = sorted(docs_path.glob("*.pdf"))
    if not pdf_files:
        return ocr_map

    print(f"\n🔎 Checking {len(pdf_files)} PDF(s) for scanned pages...")

    for pdf_path in pdf_files:
        pdf_str = str(pdf_path)

        # Check if we already have a cached OCR result for this exact file
        cache_key = hashlib.md5(
            f"{pdf_path.name}|{pdf_path.stat().st_size}|{pdf_path.stat().st_mtime}"
            .encode()
        ).hexdigest()[:12]
        ocr_out = ocr_dir / f"{pdf_path.stem}_{cache_key}_ocr.pdf"

        if ocr_out.exists():
            print(f"   ♻️  Using cached OCR for {pdf_path.name}")
            ocr_map[pdf_str] = str(ocr_out)
            continue

        print(f"   📄 Checking {pdf_path.name}...")
        if _is_scanned_pdf(pdf_str):
            print(f"   📷 Scanned PDF detected: {pdf_path.name}")
            ok = _ocr_pdf(pdf_str, str(ocr_out))
            if ok:
                ocr_map[pdf_str] = str(ocr_out)
            else:
                print(f"   ⚠️  {pdf_path.name} will be skipped (OCR unavailable/failed).")
        else:
            print(f"   ✅ Text PDF — no OCR needed: {pdf_path.name}")

    return ocr_map


# ─────────────────────────────────────────────────────────────
# FILE READERS
# ─────────────────────────────────────────────────────────────
def _read_excel(path: str) -> str:
    import pandas as pd
    try:
        ext = Path(path).suffix.lower()
        if ext == ".csv":
            return pd.read_csv(path).to_string()
        df = pd.read_excel(path, sheet_name=None)
        if isinstance(df, dict):
            return "\n\n".join(
                f"[Sheet: {s}]\n{d.to_string()}" for s, d in df.items()
            )
        return df.to_string()
    except Exception as e:
        return f"[Could not read {path}: {e}]"


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[Could not read {path}: {e}]"


def _read_docx(path: str) -> str:
    try:
        import docx2txt
        return docx2txt.process(path)
    except Exception as e:
        return f"[Could not read {path}: {e}]"


def _read_pdf(path: str) -> str:
    """
    Extract text from a PDF using pdfminer (preferred) → pypdf (fallback).
    SimpleDirectoryReader is intentionally NOT used for PDFs because it can
    fall back to raw binary reading when the PDF parser backend is missing,
    producing garbage like '%PDF-1.5 ... endobj' instead of actual text.
    """
    # ── pdfminer (most reliable for text-layer PDFs) ──────────
    try:
        from pdfminer.high_level import extract_text as _pdfminer_extract
        text = _pdfminer_extract(path) or ""
        if text.strip():
            return text
        # Empty result → fall through to pypdf
    except ImportError:
        pass
    except Exception as e:
        print(f"   ⚠️  pdfminer failed for {Path(path).name}: {e}")

    # ── pypdf fallback ────────────────────────────────────────
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        pages = []
        for page in reader.pages:
            pages.append(page.extract_text() or "")
        return "\n".join(pages)
    except ImportError:
        pass
    except Exception as e:
        print(f"   ⚠️  pypdf failed for {Path(path).name}: {e}")

    return f"[Could not extract text from {path}]"


def _read_txt(path: str) -> str:
    """Read plain-text / markdown files with UTF-8 fallback."""
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return Path(path).read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    return ""


# ─────────────────────────────────────────────────────────────
# DOCUMENT COLLECTION
# ─────────────────────────────────────────────────────────────
def _collect_documents(docs_dir: str):
    from llama_index.core import Document

    docs_path   = Path(docs_dir)
    text_docs   = []
    image_paths = []

    if not docs_path.exists():
        docs_path.mkdir(parents=True, exist_ok=True)
        print(f"📁 Created RAG docs folder: {docs_path.resolve()}")
        return [], []

    # ── OCR pre-pass: detect & fix scanned PDFs ──────────────
    # ocr_map: { original_pdf_path -> ocr_pdf_path }
    # Text PDFs are NOT in this map and are used directly.
    ocr_map = _preprocess_pdfs(docs_dir)

    for f in sorted(docs_path.rglob("*")):
        if not f.is_file():
            continue
        # Skip helper subdirectories
        if any(part in f.parts for part in ("extracted_images", "_ocr_cache")):
            continue
        ext = f.suffix.lower()
        if ext in IMAGE_EXTS:
            image_paths.append(str(f))
        elif ext in EXCEL_EXTS:
            text = _clean_text(_read_excel(str(f)))
            if _is_text_quality_ok(text):
                text_docs.append(Document(
                    text=text,
                    metadata={"file": f.name, "type": "excel"}
                ))
        elif ext in (".pptx", ".ppt"):
            text = _clean_text(_read_pptx(str(f)))
            if _is_text_quality_ok(text):
                text_docs.append(Document(
                    text=text,
                    metadata={"file": f.name, "type": "pptx"}
                ))
        elif ext in (".docx", ".doc"):
            text = _clean_text(_read_docx(str(f)))
            if _is_text_quality_ok(text):
                text_docs.append(Document(
                    text=text,
                    metadata={"file": f.name, "type": "docx"}
                ))
        elif ext == ".pdf":
            # Route through our own PDF reader — SimpleDirectoryReader can
            # silently fall back to raw binary reading when its PDF backend
            # is missing, producing '%PDF-1.5 ... endobj' garbage.
            f_str = str(f)
            pdf_to_read = ocr_map.get(f_str, f_str)  # use OCR copy if available
            if pdf_to_read != f_str:
                print(f"   📄 Using OCR'd copy: {Path(pdf_to_read).name} (was {f.name})")
            print(f"   📄 Extracting text from PDF: {f.name}")
            text = _clean_text(_read_pdf(pdf_to_read))
            if _is_text_quality_ok(text):
                text_docs.append(Document(
                    text=text,
                    metadata={"file": f.name, "type": "pdf"}
                ))
            else:
                print(f"   ⚠️  No readable text extracted from {f.name} — skipping.")
        elif ext in (".txt", ".md"):
            text = _clean_text(_read_txt(str(f)))
            if _is_text_quality_ok(text):
                text_docs.append(Document(
                    text=text,
                    metadata={"file": f.name, "type": ext.lstrip(".")}
                ))

    print(f"📄 Loaded {len(text_docs)} text docs")
    return text_docs, image_paths


# ─────────────────────────────────────────────────────────────
# MODELS
# ─────────────────────────────────────────────────────────────
def _get_embed_model():
    if EMBED_BACKEND.lower() == "huggingface":
        try:
            from llama_index.embeddings.huggingface import HuggingFaceEmbedding
            print(f"🔢 Embedding: HuggingFace ({EMBED_MODEL})")
            return HuggingFaceEmbedding(
                model_name=EMBED_MODEL,
                cache_folder="/tmp/hf_cache"
            )
        except ImportError as e:
            print(f"⚠️  HuggingFace import failed: {e}")
    from llama_index.embeddings.ollama import OllamaEmbedding
    print(f"🔢 Embedding: Ollama ({EMBED_MODEL})")
    return OllamaEmbedding(model_name=EMBED_MODEL, base_url=OLLAMA_LOCAL_URL)

def _get_llm(ollama_host: str, ollama_model: str, api_key: Optional[str] = None):
    from llama_index.llms.ollama import Ollama
    headers = {"Authorization": f"Bearer {api_key}"} if api_key else None
    return Ollama(
        model=ollama_model,
        base_url=ollama_host,
        headers=headers,
        request_timeout=60.0,
    )


def _get_chroma_vector_store():
    import chromadb
    from llama_index.vector_stores.chroma import ChromaVectorStore

    client = chromadb.PersistentClient(path=CHROMA_PERSIST_DIR)
    collection = client.get_or_create_collection(CHROMA_COLLECTION)
    return ChromaVectorStore(chroma_collection=collection)




# ─────────────────────────────────────────────────────────────
# INDEX BUILD & LOAD
# ─────────────────────────────────────────────────────────────
def build_index(
    docs_dir: str = RAG_DOCS_DIR,
    index_dir: str = RAG_INDEX_DIR,
    force_rebuild: bool = False,
    lazy: bool = False,
) -> bool:
    global _index_cache

    from llama_index.core import (
        VectorStoreIndex, StorageContext, Settings,
    )
    from llama_index.core.node_parser import SentenceSplitter

    index_path = Path(index_dir)

    if lazy:
        if not Path(CHROMA_PERSIST_DIR).exists():
            print("📚 RAG index not built yet. Say 'search my document' to build it.")
            return False
        if docs_changed(docs_dir, index_dir):
            print("📚 New/modified files detected — rebuilding index...")
            force_rebuild = True

    if force_rebuild:
        shutil.rmtree(str(CHROMA_PERSIST_DIR), ignore_errors=True)
        print("🗑️  Old Chroma data cleared.")
        _index_cache = None

    # Load from disk if available
    if Path(CHROMA_PERSIST_DIR).exists() and not force_rebuild and _index_cache is None:
        try:
            print("📂 Loading existing RAG index from Chroma...")
            Settings.embed_model = _get_embed_model()
            vector_store = _get_chroma_vector_store()
            _index_cache = VectorStoreIndex.from_vector_store(vector_store)
            print("✅ RAG index loaded.")
            return True
        except Exception as e:
            print(f"⚠️  Could not load index ({e}), rebuilding...")
            shutil.rmtree(str(CHROMA_PERSIST_DIR), ignore_errors=True)

    if _index_cache is not None and not force_rebuild:
        return True

    # ── Build fresh index ─────────────────────────────────────
    print(f"🔨 Building RAG index from '{docs_dir}'...")
    text_docs, image_paths = _collect_documents(docs_dir)

    if not text_docs and not image_paths:
        print(f"⚠️  No documents found in {docs_dir}.")
        return False


    Settings.embed_model = _get_embed_model()

    # chunk_size=512 is safe for nomic-embed-text (2048 token limit)
    splitter = SentenceSplitter(chunk_size=SAFE_CHUNK_SIZE, chunk_overlap=50)
    nodes    = splitter.get_nodes_from_documents(text_docs, show_progress=True)

    # Hard truncation: nomic-embed-text fails on inputs > ~1800 chars
    truncated = 0
    for node in nodes:
        if len(node.text) > MAX_NODE_CHARS:
            node.text = node.text[:MAX_NODE_CHARS]
            truncated += 1
    if truncated:
        print(f"✂️  Truncated {truncated} oversized chunks to {MAX_NODE_CHARS} chars.")

    print(f"📊 Total chunks: {len(nodes)} (chunk_size={SAFE_CHUNK_SIZE})")

    vector_store = _get_chroma_vector_store()
    storage = StorageContext.from_defaults(vector_store=vector_store)
    _index_cache = VectorStoreIndex(
        nodes,
        storage_context=storage,
        show_progress=True,
    )
    Path(CHROMA_PERSIST_DIR).mkdir(parents=True, exist_ok=True)
    _save_hash(docs_dir, index_dir)
    print(f"✅ Index built and saved to Chroma collection '{CHROMA_COLLECTION}'")
    return True




# ─────────────────────────────────────────────────────────────
# QUERY
# ─────────────────────────────────────────────────────────────
def query_rag(
    question: str,
    ollama_host: str,
    ollama_model: str,
    api_key: Optional[str] = None,
    docs_dir: str = RAG_DOCS_DIR,
    index_dir: str = RAG_INDEX_DIR,
):
    global _index_cache
    from llama_index.core import Settings

    if _index_cache is None:
        ok = build_index(docs_dir=docs_dir, index_dir=index_dir)
        if not ok:
            yield f"No documents found. Add files to '{docs_dir}' and say 'Boss, rebuild index'."
            return

    Settings.llm         = _get_llm(ollama_host, ollama_model, api_key)
    Settings.embed_model = _get_embed_model()

    try:
        response = _index_cache.as_query_engine(
            similarity_top_k=4,
            response_mode="compact",
            streaming=True,
        ).query(question)

        for text in response.response_gen:
            yield text

    except Exception as e:
        yield f"Sorry, I could not search your documents. Error: {e}"


def rebuild_index(docs_dir: str = RAG_DOCS_DIR, index_dir: str = RAG_INDEX_DIR) -> str:
    global _index_cache
    _index_cache = None
    ok = build_index(docs_dir=docs_dir, index_dir=index_dir, force_rebuild=True)
    return "Index rebuilt successfully." if ok else "No documents found to index."


# ─────────────────────────────────────────────────────────────
# STANDALONE TEST
# ─────────────────────────────────────────────────────────────
if __name__ == "__main__":
    env_file = Path(".env")
    if env_file.exists():
        for line in env_file.read_text().splitlines():
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            k, v = line.split("=", 1)
            os.environ.setdefault(k.strip(), v.strip().strip('"').strip("'"))

    CLOUD_HOST  = os.environ.get("OLLAMA_HOST",    "https://ollama.com")
    CLOUD_KEY   = os.environ.get("OLLAMA_API_KEY", "")
    CLOUD_MODEL = os.environ.get("OLLAMA_MODEL",   "gemma4:31b-cloud")

    print(f"🌐 Query LLM : {CLOUD_MODEL} @ {CLOUD_HOST}")
    print(f"🔢 Embed     : {EMBED_BACKEND} ({EMBED_MODEL})")
    print(f"🖼️  Vision    : {VISION_MODEL} @ {OLLAMA_LOCAL_URL}")
    print(f"📦 Chunk size: {SAFE_CHUNK_SIZE} tokens (max {MAX_NODE_CHARS} chars)")

    print("\n🔨 Building index...")
    build_index()

    print("\n💬 Ask questions about your docs (type 'quit' to exit):")
    while True:
        q = input("You: ").strip()
        if q.lower() in ("quit", "exit", "q"):
            break
        if q:
            answer = query_rag(q, CLOUD_HOST, CLOUD_MODEL, api_key=CLOUD_KEY)
            print(f"RAG: {answer}\n")